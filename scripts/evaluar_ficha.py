"""
Evalúa la veracidad de una ficha generada por Archon contra su documento fuente.

Capa 1 — Tests deterministas (siempre):
  - Integridad estructural: todas las secciones requeridas presentes
  - Metadatos: campos del documento fuente completos
  - Cobertura de términos: los términos clave aparecen en el fuente

Capa 2 — LLM-as-judge (requiere ANTHROPIC_API_KEY en el entorno):
  - Faithfulness: cada afirmación de la ficha está respaldada en el fuente
  - Completeness: no se omitieron elementos clave del fuente

Uso:
  python scripts/evaluar_ficha.py <fuente> <ficha.md>
  python scripts/evaluar_ficha.py <fuente> <ficha.md> --skip-llm
"""

import argparse
import json
import os
import re
import sys
from datetime import date
from pathlib import Path


# ── Extracción de texto del documento fuente ──────────────────────────────────

def extract_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)

def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)

def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    except ImportError:
        raise RuntimeError("python-pptx no instalado. Ejecuta: pip install python-pptx")

def extract_source_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    elif suffix in (".docx", ".doc"):
        return extract_docx(path)
    elif suffix in (".pptx", ".ppt"):
        return extract_pptx(path)
    else:
        return path.read_text(encoding="utf-8", errors="ignore")


# ── Parseo de la ficha ────────────────────────────────────────────────────────

SECCIONES = [
    "Datos del documento fuente",
    "Resumen factual",
    "Elementos prioritarios extraídos",
    "Términos clave identificados",
    "Datos relevantes adicionales",
    "Elementos excluidos",
]

METADATOS = [
    "Título original",
    "Tipo de documento",
    "Autor(es)",
    "Fecha",
    "Fuente / origen",
]

def parse_ficha(path: Path) -> dict:
    text = path.read_text(encoding="utf-8")
    result = {"titulo": "", "secciones": {}, "metadatos": {}, "terminos": [], "raw": text}

    m = re.search(r"^#\s+(.+)$", text, re.MULTILINE)
    result["titulo"] = m.group(1).strip() if m else path.stem

    for sec in SECCIONES:
        m = re.search(rf"## {re.escape(sec)}\n(.*?)(?=\n## |\Z)", text, re.DOTALL)
        result["secciones"][sec] = m.group(1).strip() if m else None

    datos_m = re.search(r"## Datos del documento fuente\n(.*?)(?=\n## |\Z)", text, re.DOTALL)
    if datos_m:
        for campo in METADATOS:
            m = re.search(rf"- {re.escape(campo)}:\s*(.+)", datos_m.group(1))
            result["metadatos"][campo] = m.group(1).strip() if m else ""

    tc_m = re.search(r"## Términos clave identificados\n(.*?)(?=\n## |\Z)", text, re.DOTALL)
    if tc_m:
        bullet_terms = re.findall(r"^[-•]\s+(.+?)(?:\s*/.*|\s*\(.*)?$", tc_m.group(1), re.MULTILINE)
        result["terminos"] = [t.strip() for t in bullet_terms if t.strip()][:25]

    return result


# ── Capa 1: tests deterministas ───────────────────────────────────────────────

VACIO_PLACEHOLDER = "[No disponible en el documento fuente]"

def test_estructura(ficha: dict) -> dict:
    out = {}
    for sec in SECCIONES:
        content = ficha["secciones"].get(sec)
        if content is None:
            out[sec] = "AUSENTE"
        elif VACIO_PLACEHOLDER in content:
            out[sec] = "NO_DISPONIBLE"
        elif len(content.strip()) < 10:
            out[sec] = "VACIA"
        else:
            out[sec] = "PRESENTE"
    return out

def test_metadatos(ficha: dict) -> dict:
    out = {}
    for campo in METADATOS:
        val = ficha["metadatos"].get(campo, "")
        out[campo] = "PRESENTE" if val and val not in ("[No disponible]", "") else "VACIO"
    return out

def test_terminos(ficha: dict, source_text: str) -> dict:
    source_lower = source_text.lower()
    out = {}
    for term in ficha["terminos"]:
        clean = re.sub(r"\(.*?\)", "", term).strip().lower()
        words = [w for w in clean.split() if len(w) > 4]
        if not words:
            words = [clean]
        out[term] = "ENCONTRADO" if any(w in source_lower for w in words) else "NO_ENCONTRADO"
    return out


# ── Capa 2: LLM-as-judge ─────────────────────────────────────────────────────

FAITH_PROMPT = """\
Eres un evaluador de calidad de fichas documentales. Verifica si las afirmaciones \
de la sección son fieles al documento fuente.

DOCUMENTO FUENTE:
{source}

SECCIÓN — {section}:
{content}

Para cada afirmación factual relevante determina:
- SOPORTADA: encontrada directamente en el fuente
- INFERIDA: implícita pero no explícita
- NO_ENCONTRADA: sin evidencia en el fuente

Responde SOLO con JSON:
{{
  "faithfulness_score": <0-100>,
  "claims_evaluated": <n>,
  "supported": <n>,
  "inferred": <n>,
  "not_found": <n>,
  "issues": [{{"claim": "...", "status": "INFERIDA|NO_ENCONTRADA", "note": "..."}}]
}}"""

COMPLETENESS_PROMPT = """\
Eres un evaluador de calidad de fichas documentales. Identifica información \
importante del fuente que NO está capturada en la ficha.

DOCUMENTO FUENTE:
{source}

FICHA COMPLETA:
{ficha}

Responde SOLO con JSON:
{{
  "completeness_score": <0-100>,
  "key_elements_in_source": <n>,
  "captured": <n>,
  "omissions": [{{"element": "...", "importance": "ALTA|MEDIA"}}]
}}"""

def truncate(text: str, max_chars: int = 60000) -> str:
    if len(text) <= max_chars:
        return text
    half = max_chars // 2
    return text[:half] + "\n\n[... truncado ...]\n\n" + text[-half:]

def parse_json_response(text: str) -> dict:
    m = re.search(r"\{.*\}", text, re.DOTALL)
    if m:
        return json.loads(m.group())
    raise ValueError("No se encontró JSON en la respuesta")

def run_llm_eval(ficha: dict, ficha_path: Path, source_text: str) -> dict:
    try:
        import anthropic
    except ImportError:
        return {"error": "anthropic no instalado — ejecuta: pip install anthropic"}

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        return {"error": "ANTHROPIC_API_KEY no está configurada en el entorno"}

    client = anthropic.Anthropic(api_key=api_key)
    source = truncate(source_text)
    results = {"secciones": {}, "completeness": None}

    for sec in ["Resumen factual", "Elementos prioritarios extraídos", "Datos relevantes adicionales"]:
        content = ficha["secciones"].get(sec, "")
        if not content or len(content) < 20:
            continue
        try:
            prompt = FAITH_PROMPT.format(
                source=source,
                section=sec,
                content=content[:4000],
            )
            resp = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=1024,
                messages=[{"role": "user", "content": prompt}],
            )
            results["secciones"][sec] = parse_json_response(resp.content[0].text)
        except Exception as e:
            results["secciones"][sec] = {"error": str(e)}

    try:
        prompt = COMPLETENESS_PROMPT.format(
            source=source,
            ficha=ficha["raw"][:5000],
        )
        resp = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1024,
            messages=[{"role": "user", "content": prompt}],
        )
        results["completeness"] = parse_json_response(resp.content[0].text)
    except Exception as e:
        results["completeness"] = {"error": str(e)}

    return results


# ── Generación del reporte .md ────────────────────────────────────────────────

def fmt_score(score) -> str:
    if isinstance(score, int):
        bar = "█" * (score // 10) + "░" * (10 - score // 10)
        return f"{score}% {bar}"
    return str(score)

def build_report(ficha, ficha_path, source_path, estructura, metadatos, terminos, llm=None) -> str:
    lines = []

    lines += [
        f"# Evaluación — {ficha['titulo']}",
        "",
        f"**Fuente:** `{source_path.name}`  ",
        f"**Ficha:** `{ficha_path.name}`  ",
        f"**Fecha:** {date.today().isoformat()}  ",
        f"**Archon:** v1.03",
        "",
        "---",
        "",
        "## Capa 1 — Tests deterministas",
        "",
    ]

    # Estructura
    struct_ok = sum(1 for v in estructura.values() if v == "PRESENTE")
    struct_score = round(struct_ok / len(estructura) * 100)
    lines += ["### Integridad estructural", ""]
    lines += ["| Sección | Estado |", "|---|---|"]
    for sec, status in estructura.items():
        prefix = "✓" if status == "PRESENTE" else "✗"
        lines.append(f"| {sec} | {prefix} {status} |")
    lines += ["", f"**Score: {fmt_score(struct_score)}**", ""]

    # Metadatos
    meta_ok = sum(1 for v in metadatos.values() if v == "PRESENTE")
    meta_score = round(meta_ok / len(metadatos) * 100) if metadatos else 0
    lines += ["### Metadatos", ""]
    lines += ["| Campo | Estado |", "|---|---|"]
    for campo, status in metadatos.items():
        prefix = "✓" if status == "PRESENTE" else "✗"
        lines.append(f"| {campo} | {prefix} {status} |")
    lines += ["", f"**Score: {fmt_score(meta_score)}**", ""]

    # Términos
    lines += ["### Cobertura de términos clave", ""]
    if terminos:
        found = sum(1 for v in terminos.values() if v == "ENCONTRADO")
        term_score = round(found / len(terminos) * 100)
        for term, status in terminos.items():
            prefix = "+" if status == "ENCONTRADO" else "−"
            lines.append(f"- [{prefix}] {term}")
        lines += ["", f"**Score: {fmt_score(term_score)}**", ""]
    else:
        term_score = None
        lines += ["_No se extrajeron términos clave de la ficha._", ""]

    # Capa 2
    faith_scores = []
    comp_score = None

    if llm:
        lines += ["---", "", "## Capa 2 — LLM-as-judge", ""]

        if "error" in llm:
            lines += [f"> **No disponible:** {llm['error']}", ""]
        else:
            for sec, data in llm.get("secciones", {}).items():
                lines += [f"### Faithfulness — {sec}", ""]
                if "error" in data:
                    lines.append(f"_Error: {data['error']}_")
                else:
                    sc = data.get("faithfulness_score", 0)
                    faith_scores.append(sc)
                    lines.append(f"**Score: {fmt_score(sc)}**  ")
                    lines.append(
                        f"Afirmaciones: {data.get('claims_evaluated','?')} — "
                        f"Soportadas: {data.get('supported','?')} | "
                        f"Inferidas: {data.get('inferred','?')} | "
                        f"No encontradas: {data.get('not_found','?')}"
                    )
                    issues = data.get("issues", [])
                    if issues:
                        lines += ["", "**Issues:**"]
                        for i in issues:
                            lines.append(f"- `{i.get('status','')}` — _{i.get('claim','')}_")
                            lines.append(f"  {i.get('note','')}")
                lines.append("")

            comp = llm.get("completeness")
            if comp and "error" not in comp:
                comp_score = comp.get("completeness_score", 0)
                lines += ["### Completeness", ""]
                lines.append(f"**Score: {fmt_score(comp_score)}**  ")
                lines.append(
                    f"Elementos en fuente: {comp.get('key_elements_in_source','?')} — "
                    f"Capturados: {comp.get('captured','?')}"
                )
                omissions = comp.get("omissions", [])
                if omissions:
                    lines += ["", "**Posibles omisiones:**"]
                    for o in omissions:
                        lines.append(f"- [{o.get('importance','')}] {o.get('element','')}")
                lines.append("")

    # Resumen
    lines += ["---", "", "## Resumen", ""]
    all_scores = [struct_score, meta_score]
    if term_score is not None:
        all_scores.append(term_score)

    lines += ["| Dimensión | Score |", "|---|---|"]
    lines.append(f"| Integridad estructural | {struct_score}% |")
    lines.append(f"| Metadatos | {meta_score}% |")
    if term_score is not None:
        lines.append(f"| Cobertura de términos | {term_score}% |")

    if faith_scores:
        avg_faith = round(sum(faith_scores) / len(faith_scores))
        all_scores.append(avg_faith)
        lines.append(f"| Faithfulness (LLM) | {avg_faith}% |")
    if comp_score is not None:
        all_scores.append(comp_score)
        lines.append(f"| Completeness (LLM) | {comp_score}% |")

    global_score = round(sum(all_scores) / len(all_scores))
    lines.append(f"| **SCORE GLOBAL** | **{global_score}%** |")

    if not llm or "error" in llm:
        lines += [
            "",
            "_Para evaluación semántica completa, configura ANTHROPIC_API_KEY y ejecuta sin `--skip-llm`._",
        ]

    return "\n".join(lines) + "\n"


# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Evalúa la veracidad de una ficha de Archon contra su documento fuente."
    )
    parser.add_argument("fuente", help="Documento fuente (PDF, DOCX, PPTX)")
    parser.add_argument("ficha",  help="Ficha .md generada por Archon")
    parser.add_argument("--skip-llm", action="store_true",
                        help="Solo ejecuta tests deterministas (Capa 1)")
    args = parser.parse_args()

    source_path = Path(args.fuente)
    ficha_path  = Path(args.ficha)

    for p in (source_path, ficha_path):
        if not p.exists():
            print(f"Error: no se encontró {p}")
            sys.exit(1)

    print("Extrayendo texto del fuente...")
    source_text = extract_source_text(source_path)
    if len(source_text.strip()) < 100:
        print("Advertencia: poco texto extraído — puede ser un PDF escaneado.")

    print("Parseando ficha...")
    ficha = parse_ficha(ficha_path)

    print("Capa 1 — tests deterministas...")
    estructura = test_estructura(ficha)
    metadatos  = test_metadatos(ficha)
    terminos   = test_terminos(ficha, source_text)

    llm = None
    if not args.skip_llm:
        print("Capa 2 — LLM-as-judge...")
        llm = run_llm_eval(ficha, ficha_path, source_text)

    print("Generando reporte...")
    report = build_report(ficha, ficha_path, source_path,
                          estructura, metadatos, terminos, llm)

    base       = Path(__file__).parent.parent
    output_dir = base / "salidas" / "evaluaciones"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"Evaluacion — {ficha_path.stem}.md"
    output_path.write_text(report, encoding="utf-8")
    print(f"Reporte guardado en: {output_path}")


if __name__ == "__main__":
    main()
